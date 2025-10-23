 # Стандартные библиотеки Python
import csv
import http.client
import io
import logging
import re
from datetime import datetime
from urllib.parse import urlparse

# Сторонние библиотеки
import requests
import telebot
from bs4 import BeautifulSoup
from openai import OpenAI
from PyPDF2 import PdfReader
from docx import Document  # .docx
from pptx import Presentation  # .pptx
import openpyxl  # .xlsx
import chardet  # для определения кодировки .txt/.csv

# Google Colab (опционально)
try:
    from google.colab import userdata
    import os
except ImportError:
    # Для локального запуска без Google Colab
    import os
    userdata = type('UserData', (), {
        'get': lambda self, key: os.environ.get(key)
    })()

# Загрузка .env для локального запуска
from bot_env_loader import *

# Настройка логирования
logging.basicConfig(level=logging.INFO)

# Получение токена Telegram API из переменной окружения
# Получение токена Telegram API из userdata
API_TOKEN = userdata.get('TELEGRAM_BOT_TOKEN')
if not API_TOKEN:
    logging.critical("Необходимо установить переменную окружения TELEGRAM_BOT_TOKEN с токеном вашего бота!")
    exit(1)  # Завершаем программу, если токен не найден

bot = telebot.TeleBot(API_TOKEN)

# Получение API-ключей OpenAI из переменных окружения
OPENAI_API_KEY = userdata.get("OPENAI_API_KEY")
BACKUP_OPENAI_API_KEY = userdata.get("BACKUP_OPENAI_API_KEY")

if not OPENAI_API_KEY:
    logging.critical("Необходимо установить переменную окружения OPENAI_API_KEY с вашим API-ключом!")
    exit(1)  # Завершаем программу, если API-ключ не найден

# Настройка OpenAI с основным ключом
client = OpenAI(
    api_key=OPENAI_API_KEY,
    base_url="https://api.proxyapi.ru/openai/v1",
)

# Настройка резервного клиента (если есть альтернативный ключ)
backup_client = None
if BACKUP_OPENAI_API_KEY:
    backup_client = OpenAI(
        api_key=BACKUP_OPENAI_API_KEY,
        base_url="https://api.proxyapi.ru/openai/v1",
    )

# Путь к файлу логов (для локального запуска — файл будет создан в рабочей директории)
file_path = 'telegram_bot_logs.csv'  # Локальный путь для логирования сообщений

# Словарь для хранения истории разговора
conversation_history = {}

def process_url_in_text(text, bot, chat_id):
    """
    Ищет URL в тексте и, если находит, извлекает текст с веб-страницы.

    Args:
        text (str): Текст для поиска URL.
        bot (telebot.TeleBot): Экземпляр бота.
        chat_id (int): ID чата.

    Returns:
        str: Объединенный текст (исходный текст + текст с веб-страницы) или исходный текст, если URL не найден.
    """
    url_match = re.search(r'(http[s]?://[^\s]+)', text)
    if url_match:
        url = url_match.group(0)  # Первая найденная ссылка
        extracted_text = extract_text_from_url(url)

        if extracted_text:
            return f"{text}\n\n{extracted_text}"
        else:
            bot.reply_to(chat_id, "Не удалось извлечь текст из ссылки.")
            return text
    else:
        return text

# Функция для извлечения текста из URL
import random

def extract_text_from_url(url):
    # Расширенный список User-Agent, выбор случайного при каждом запросе
    USER_AGENTS = [
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:119.0) Gecko/20100101 Firefox/119.0",
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/118.0.5993.88 Safari/537.36",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 11_5_2) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/117.0.0.0 Safari/537.36",
        "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/116.0.0.0 Safari/537.36"
    ]
    headers = {
        "User-Agent": random.choice(USER_AGENTS),
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8",
        "Accept-Language": "ru-RU,ru;q=0.9,en-US;q=0.8,en;q=0.7",
        "Referer": url,
        "Connection": "keep-alive",
        "DNT": "1",
        # "Cookie": ""  # Можно добавить пустую строку или свой набор,
    }
    try:
        resp = requests.get(url, headers=headers, timeout=10, allow_redirects=True)
        if resp.status_code != 200:
            return f"Ошибка: {resp.status_code}"
        # Пытаемся определить правильную кодировку
        resp.encoding = resp.apparent_encoding if resp.encoding is None else resp.encoding
        page_content = resp.text
        soup = BeautifulSoup(page_content, 'html.parser')
        text_content = soup.get_text(separator='\n')
        cleaned_text = "\n".join(line.strip() for line in text_content.splitlines() if line.strip())
        # Опционально: урезать по длине, если слишком много текста
        if len(cleaned_text) > 5000:
            cleaned_text = cleaned_text[:5000] + "..."
        return cleaned_text
    except Exception as e:
        return f"Произошла ошибка: {e}"

# Команда /start
@bot.message_handler(commands=['start'])
def send_welcome(message):
    chat_id = message.chat.id
    # Инициализация истории разговора для нового чата
    if chat_id not in conversation_history:
        conversation_history[chat_id] = []
    bot.reply_to(message, "Добро пожаловать в канал 'Это не канал'! Как я могу помочь? Бот версии 21_01_2025 г")

# Обработка текстовых сообщений
@bot.message_handler(content_types=['text'])
def handle_text_message(message):
    chat_id = message.chat.id
    user_message = message.text  # Текст сообщения пользователя
    message_type = 'text'  # Указываем тип сообщения как текстовое

    # Проверяем, содержит ли сообщение текст "http"
    if "http" in user_message:
        # Сохраняем исходный текст сообщения
        original_message = user_message

        # Извлекаем первую ссылку из текста
        import re
        url_match = re.search(r'(http[s]?://[^\s]+)', user_message)
        if url_match:
            url = url_match.group(0)  # Первая найденная ссылка

            # Пытаемся извлечь текст с веб-страницы
            extracted_text = extract_text_from_url(url)

            if extracted_text:
                # Объединяем текст сообщения с извлеченным текстом
                user_message = f"{original_message}\n\n{extracted_text}"
                process_message(message, user_message, message_type, chat_id)
            else:
                bot.reply_to(message, "Не удалось извлечь текст из ссылки.")
        else:
            bot.reply_to(message, "Ссылка не найдена в сообщении.")
    else:
        # Если сообщение не содержит "http", обрабатываем его как обычный текст
        process_message(message, user_message, message_type, chat_id)

@bot.message_handler(content_types=['photo'])
def handle_photo_message(message):
    chat_id = message.chat.id
    user_message = message.caption if message.caption else "Фото без подписи"
    message_type = 'photo'

    # Обрабатываем URL в подписи, если он есть
    user_message += process_url_in_text(user_message, bot, chat_id)

    try:
        # Получаем file_id самой большой версии фотографии
        file_id = message.photo[-1].file_id

        # Получаем информацию о файле
        file_info = bot.get_file(file_id)
        file_path = file_info.file_path
        image_url = f'https://api.telegram.org/file/bot{API_TOKEN}/{file_path}'

        logging.info(f"URL изображения: {image_url}")  # Логируем URL

    except Exception as e:
        logging.error(f"Ошибка при получении URL изображения из Telegram: {e}")
        user_message += "\nНе удалось получить URL изображения."
        process_message(message, user_message, message_type, chat_id)
        return  # Выходим из функции, чтобы избежать дальнейших ошибок

    try:
        # Запрашиваем описание изображения у OpenAI
        response = make_openai_request(
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Что на этом изображении? Дай краткое описание на русском языке."},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": image_url,
                            },
                        },
                    ],
                }
            ],
            model="gpt-4o-mini",
            max_tokens=300,
        )

        logging.info(f"Ответ от OpenAI Vision API: {response}")  # Логируем полный ответ

        # Извлекаем описание изображения из ответа OpenAI
        image_description = response.choices[0].message.content

        # Добавляем описание изображения к сообщению пользователя
        user_message += f"\nОписание изображения: {image_description}"

    except Exception as e:
        logging.error(f"Ошибка при обращении к OpenAI Vision API: {e}")
        user_message += "\nНе удалось получить описание изображения."

    process_message(message, user_message, message_type, chat_id)


@bot.message_handler(content_types=['document'])
def handle_document_message(message):
    chat_id = message.chat.id
    user_message = message.caption if message.caption else "Документ без подписи"
    message_type = 'document'

    file_name = message.document.file_name
    file_ext = (file_name.rsplit('.', 1)[-1] if '.' in file_name else '').lower()

    # Обрабатываем URL в подписи, если он есть
    user_message += process_url_in_text(user_message, bot, chat_id)

    try:
        # Получаем информацию о файле
        file_info = bot.get_file(message.document.file_id)
        file_path = file_info.file_path
        doc_url = f'https://api.telegram.org/file/bot{API_TOKEN}/{file_path}'

        logging.info(f"URL документа: {doc_url}")  # Логируем URL

    except Exception as e:
        logging.error(f"Ошибка при получении URL документа из Telegram: {e}")
        user_message += "\nНе удалось получить URL документа."
        process_message(message, user_message, message_type, chat_id)
        return

    try:
        # Скачиваем файл
        response = requests.get(doc_url)
        response.raise_for_status()
        # Универсальное извлечение текста для разных форматов
        file_bytes = response.content
        extracted_text = extract_text_from_document(file_bytes, file_name)

        if not extracted_text.strip():
            user_message += "\nНе удалось извлечь текст из документа."
            process_message(message, user_message, message_type, chat_id)
            return

        # Ограничиваем размер текста для API (примерно 4000 токенов)
        if len(extracted_text) > 12000:  # Примерно 4000 токенов
            extracted_text = extracted_text[:12000] + "\n... (текст обрезан из-за ограничений)"

        logging.info(f"Извлеченный текст из документа: {extracted_text[:500]}...")  # Логируем начало текста

    except Exception as e:
        logging.error(f"Ошибка при обработке документа: {e}")
        user_message += "\nНе удалось обработать документ."
        process_message(message, user_message, message_type, chat_id)
        return

    try:
        # Запрашиваем анализ документа у OpenAI
        response = make_openai_request(
            messages=[
                {
                    "role": "user",
                    "content": f"""Проанализируй этот документ и дай краткое описание на русском языке.

                    Включи в описание:
                    - Тип документа
                    - Основную тему/содержание
                    - Ключевые пункты
                    - Количество страниц/слайдов/листов (если видно из текста)

                    Текст документа:
                    {extracted_text}"""
                }
            ],
            model="gpt-4o-mini",
            max_tokens=500,
        )

        logging.info(f"Ответ от OpenAI для документа: {response}")  # Логируем полный ответ

        # Извлекаем анализ документа из ответа OpenAI
        doc_analysis = response.choices[0].message.content

        # Добавляем анализ к сообщению пользователя
        user_message += f"\n\нАнализ документа:\n{doc_analysis}"

    except Exception as e:
        logging.error(f"Ошибка при обращении к OpenAI для анализа документа: {e}")
        user_message += "\nНе удалось получить анализ документа."

    process_message(message, user_message, message_type, chat_id)

def extract_text_from_document(file_bytes: bytes, file_name: str) -> str:
    """
    Извлекает текст из поддерживаемых офисных форматов: pdf, txt, docx, pptx, xlsx, csv.
    Для .doc (старый формат) извлечение не поддерживается надёжно — будет возвращено пусто.
    """
    ext = (file_name.rsplit('.', 1)[-1] if '.' in file_name else '').lower()
    data = io.BytesIO(file_bytes)

    try:
        if ext == 'pdf':
            reader = PdfReader(data)
            parts = []
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    text = page.extract_text() or ''
                    if text.strip():
                        parts.append(f"\n--- Страница {page_num} ---\n{text}")
                except Exception as e:
                    logging.warning(f"Не удалось извлечь текст со страницы {page_num}: {e}")
                    continue
            return ''.join(parts)

        if ext in ('txt', 'log', 'md', 'csv'):  # простые текстовые
            detection = chardet.detect(file_bytes)
            encoding = detection.get('encoding') or 'utf-8'
            try:
                return file_bytes.decode(encoding, errors='ignore')
            except Exception:
                return file_bytes.decode('utf-8', errors='ignore')

        if ext == 'docx':
            doc = Document(data)
            return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())

        if ext == 'pptx':
            presentation = Presentation(data)
            texts = []
            for slide_idx, slide in enumerate(presentation.slides, 1):
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        slide_text_parts.append(shape.text)
                if slide_text_parts:
                    texts.append(f"\n--- Слайд {slide_idx} ---\n" + '\n'.join(slide_text_parts))
            return '\n'.join(texts)

        if ext in ('xlsx', 'xlsm', 'xltx'):
            wb = openpyxl.load_workbook(data, data_only=True, read_only=True)
            texts = []
            for sheet in wb.worksheets:
                texts.append(f"\n--- Лист: {sheet.title} ---")
                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    row_count += 1
                    # Ограничим вывод первыми 1000 строками на лист
                    if row_count > 1000:
                        texts.append("... (лист обрезан по ограничению строк)")
                        break
                    row_text = '\t'.join('' if v is None else str(v) for v in row)
                    if row_text.strip():
                        texts.append(row_text)
            return '\n'.join(texts)

        # .doc (старый формат) и прочие не поддерживаемые
        logging.warning(f"Формат .{ext} пока не поддерживается для извлечения текста")
        return ''
    except Exception as e:
        logging.error(f"Ошибка извлечения текста из документа ({file_name}): {e}")
        return ''

# Обработка сообщений с видео
@bot.message_handler(content_types=['video'])
def handle_video_message(message):
    chat_id = message.chat.id
    user_message = message.caption if message.caption else "Видео без подписи"
    message_type = 'video'

    # Обрабатываем URL в подписи, если он есть
    user_message += process_url_in_text(user_message, bot, chat_id)

    process_message(message, user_message, message_type, chat_id)


# Обработка голосовых сообщений
@bot.message_handler(content_types=['voice'])
def handle_voice_message(message):
    chat_id = message.chat.id
    message_type = 'voice'

    try:
        # Получаем информацию о файле
        file_id = message.voice.file_id
        file_info = bot.get_file(file_id)
        file_path = file_info.file_path
        file_url = f'https://api.telegram.org/file/bot{API_TOKEN}/{file_path}'

        # Скачиваем файл
        response = requests.get(file_url)
        response.raise_for_status()  # Проверяем на ошибки

        # Открываем файл как бинарный
        audio_file = io.BytesIO(response.content)

        # Транскрибируем аудио
        transcription = make_whisper_request(
            file_data=audio_file,
            filename=file_path
        )

        # Получаем транскрибированный текст
        transcribed_text = transcription.text

        # Формируем сообщение пользователю: сначала подпись, потом транскрипция
        user_message = message.caption if message.caption else ""  # Получаем подпись
        user_message = process_url_in_text(user_message, bot, chat_id)  # Обрабатываем URL в подписи, если он есть
        user_message += f"\nТранскрипция аудио: {transcribed_text}"  # Добавляем транскрипцию

    except Exception as e:
        logging.error(f"Ошибка при транскрибации аудио: {e}")
        bot.reply_to(message, "Произошла ошибка при транскрибации аудио.")

    process_message(message, user_message, message_type, chat_id)


# Обработка аудио сообщений
@bot.message_handler(content_types=['audio'])
def handle_audio_message(message):
    chat_id = message.chat.id
    message_type = 'audio'

    try:
        # Получаем информацию о файле
        file_id = message.audio.file_id  # Изменено с message.voice.file_id на message.audio.file_id
        file_info = bot.get_file(file_id)
        file_path = file_info.file_path
        file_url = f'https://api.telegram.org/file/bot{API_TOKEN}/{file_path}'

        # Скачиваем файл
        response = requests.get(file_url)
        response.raise_for_status()  # Проверяем на ошибки

        # Открываем файл как бинарный
        audio_file = io.BytesIO(response.content)

        # Транскрибируем аудио
        transcription = make_whisper_request(
            file_data=audio_file,
            filename=file_path
        )

        # Получаем транскрибированный текст
        transcribed_text = transcription.text

         # Формируем сообщение пользователю: сначала подпись, потом транскрипция
        user_message = message.caption if message.caption else ""  # Получаем подпись
        user_message = process_url_in_text(user_message, bot, chat_id) # Обрабатываем URL в подписи, если он есть
        user_message += f"\nТранскрипция аудио: {transcribed_text}"  # Добавляем транскрипцию

    except Exception as e:
        logging.error(f"Ошибка при транскрибации аудио: {e}")
        bot.reply_to(message, "Произошла ошибка при транскрибации аудио.")

    process_message(message, user_message, message_type, chat_id)

# Обработка опросов
@bot.message_handler(content_types=['poll'])
def handle_poll_message(message):
    chat_id = message.chat.id
    user_message = f"Опрос: {message.poll.question}"
    message_type = 'poll'
    process_message(message, user_message, message_type, chat_id)

# Функция для инициализации файла (создаем файл и пишем заголовки, если он не существует)
def initialize_log_file():
    if not os.path.exists(file_path):
        with open(file_path, 'w', newline='', encoding='utf-8') as file:
            writer = csv.writer(file)
            writer.writerow(['chat_id', 'datetime', 'message', 'message_type', 'ai_response'])

# Функция для записи данных в файл
def log_to_file(chat_id, user_message, message_type, ai_response):
    current_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')  # Текущее время
    with open(file_path, 'a', newline='', encoding='utf-8') as file:
        writer = csv.writer(file)
        writer.writerow([chat_id, current_time, user_message, message_type, ai_response])

# Функция для выполнения запроса к Whisper API с fallback логикой
def make_whisper_request(file_data, filename):
    """
    Выполняет запрос к Whisper API с автоматическим переключением на резервный ключ при ошибке 402.
    
    Args:
        file_data: Данные аудиофайла
        filename: Имя файла
    
    Returns:
        Транскрипция или None при ошибке
    """
    # Сначала пробуем основной ключ
    try:
        transcription = client.audio.transcriptions.create(
            model="whisper-1",
            file=(filename, file_data)
        )
        return transcription
    except Exception as e:
        error_str = str(e)
        # Проверяем, является ли это ошибкой 402 (недостаточно средств)
        if "402" in error_str or "Payment Required" in error_str or "Insufficient balance" in error_str:
            logging.warning(f"Основной API ключ исчерпан (402), переключаемся на резервный для Whisper: {e}")
            
            # Если есть резервный клиент, пробуем его
            if backup_client:
                try:
                    transcription = backup_client.audio.transcriptions.create(
                        model="whisper-1",
                        file=(filename, file_data)
                    )
                    logging.info("Успешно использован резервный API ключ для Whisper")
                    return transcription
                except Exception as backup_error:
                    logging.error(f"Ошибка и в резервном ключе для Whisper: {backup_error}")
                    raise backup_error
            else:
                logging.error("Резервный API ключ не настроен для Whisper")
                raise e
        else:
            # Если это не ошибка 402, просто пробрасываем исключение
            raise e

# Функция для выполнения запроса к OpenAI с fallback логикой
def make_openai_request(messages, model="gpt-3.5-turbo-1106", max_tokens=None):
    """
    Выполняет запрос к OpenAI с автоматическим переключением на резервный ключ при ошибке 402.
    
    Args:
        messages: Список сообщений для OpenAI
        model: Модель для использования
        max_tokens: Максимальное количество токенов
    
    Returns:
        Ответ от OpenAI или None при ошибке
    """
    # Сначала пробуем основной ключ
    try:
        params = {
            "model": model,
            "messages": messages
        }
        if max_tokens:
            params["max_tokens"] = max_tokens
            
        response = client.chat.completions.create(**params)
        return response
    except Exception as e:
        error_str = str(e)
        # Проверяем, является ли это ошибкой 402 (недостаточно средств)
        if "402" in error_str or "Payment Required" in error_str or "Insufficient balance" in error_str:
            logging.warning(f"Основной API ключ исчерпан (402), переключаемся на резервный: {e}")
            
            # Если есть резервный клиент, пробуем его
            if backup_client:
                try:
                    response = backup_client.chat.completions.create(**params)
                    logging.info("Успешно использован резервный API ключ")
                    return response
                except Exception as backup_error:
                    logging.error(f"Ошибка и в резервном ключе: {backup_error}")
                    raise backup_error
            else:
                logging.error("Резервный API ключ не настроен")
                raise e
        else:
            # Если это не ошибка 402, просто пробрасываем исключение
            raise e

# Общая функция для обработки сообщений
def process_message(message, user_message, message_type, chat_id):
    # Проверяем, существует ли история для данного chat_id
    if chat_id not in conversation_history:
        conversation_history[chat_id] = []

    # Добавляем сообщение пользователя в историю разговора
    conversation_history[chat_id].append({"role": "user", "content": user_message})
    logging.info(f"Получено сообщение от пользователя: {user_message} (Тип: {message_type})")
    
    # Запрос к OpenAI с историей разговора
    try:
        chat_completion = make_openai_request(
            messages=conversation_history[chat_id] + [
                {"role": "system", "content": (
                    "Вы бот-администратор в телеграм-канале 'Это не канал'. Ваша задача — пересказывать на русском языке "
                    "подписчикам материалы, присылаемые в канал. Формируйте краткий (не более 3000 знаков) и интересный пересказ, "
                    "ориентируясь на следующие принципы:\n\n"
                    "1. Прочитайте пост. Если в посте указана ссылка, предположите, что она содержит дополнительную информацию. "
                    "Попробуйте дать пересказ, основываясь на теме, изложенной в посте, а также возможных контекстах.\n"
                    "2. Пересказ оформляйте структурно:\n"
                    "- Введение: кратко объясните, о чем материал и почему он важен.\n"
                    "- Основная часть: изложите ключевые моменты материала простым языком, подчеркивая суть. Разделяйте текст на абзацы.\n"
                    "- Заключение: сделайте выводы, предложите рекомендации или задайте вопрос для вовлечения подписчиков.\n"
                    "3. Если пост содержит только ссылку, составьте предположительный пересказ на основе общего контекста и доступной информации. "
                    "Укажите, что пересказ основан на интерпретации.\n"
                    "4. Указывайте источник информации в конце текста (например: 'Источник: ссылка из поста').\n\n"
                    "Общайтесь с читателями вежливо, от мужского лица, используя 'Вы'.\n\n"
                    "Включайте эмодзи для акцентирования ключевых моментов, таких как:\n"
                    "- 🔍 для выделения важных деталей,\n"
                    "- 📌 для ключевых тезисов,\n"
                    "- 🌟 для рекомендаций.\n\n"
                    "Следите за тем, чтобы текст был легко читаем на русском языке и не перегружен эмодзи. Старайтесь создавать увлекательные посты, чтобы подписчики захотели прочитать оригинал."
                )},
                {"role": "user", "content": user_message}
            ]
        )
        # Получаем ответ от AI
        ai_response = chat_completion.choices[0].message.content
        bot.reply_to(message, ai_response)

        # Логирование данных в файл
        log_to_file(chat_id, user_message, message_type, ai_response)

        # Добавляем ответ AI в историю разговора
        conversation_history[chat_id].append({"role": "assistant", "content": ai_response})
    except Exception as e:
        logging.error(f"Ошибка при обращении к OpenAI: {e}")
        bot.reply_to(message, "Извините, произошла ошибка при обработке вашего запроса.")

# Запуск бота
if __name__ == '__main__':
    initialize_log_file()  # Инициализация файла логов
    logging.info(f"Бот запущен. Лог-файл: {os.path.abspath(file_path)}")
    bot.polling(none_stop=True)
